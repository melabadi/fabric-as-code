#!/usr/bin/env python
"""
build_pptx.py — generate the Fabric-as-Code slide deck (PowerPoint).

Usage:
    python docs/build_pptx.py

Produces: docs/Fabric-as-Code.pptx (embeds the screenshots in docs/screenshots).
Requires: python-pptx, Pillow  (pip install python-pptx Pillow)
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = Path(__file__).resolve().parent
SHOTS = HERE / "screenshots"
OUT = HERE / "Fabric-as-Code.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

FABRIC = RGBColor(0x11, 0x73, 0x65)
FABRIC_LT = RGBColor(0xE3, 0xF1, 0xEE)
DARK = RGBColor(0x20, 0x20, 0x20)
GREY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def band(slide, height=Inches(1.0), color=FABRIC):
    box = slide.shapes.add_shape(1, 0, 0, SW, height)
    box.fill.solid(); box.fill.fore_color.rgb = color
    box.line.fill.background(); box.shadow.inherit = False
    return box


def textbox(slide, left, top, width, height, text, size=18, color=DARK,
            bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return tb


def header(slide, title):
    band(slide)
    textbox(slide, Inches(0.6), Inches(0.12), Inches(12.1), Inches(0.76), title,
            size=28, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)


def title_slide(title, subtitle, tagline):
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = FABRIC; bg.line.fill.background()
    bg.shadow.inherit = False
    textbox(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.5), title,
            size=54, color=WHITE, bold=True)
    textbox(s, Inches(0.85), Inches(3.9), Inches(11.6), Inches(1.2), subtitle,
            size=24, color=WHITE)
    textbox(s, Inches(0.85), Inches(5.2), Inches(11.6), Inches(0.6), tagline,
            size=16, color=RGBColor(0xD7, 0xEF, 0xE9))


def content_slide(title, bullets):
    s = prs.slides.add_slide(BLANK)
    header(s, title)
    body = s.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.6))
    tf = body.text_frame; tf.word_wrap = True
    for i, (txt, lvl) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        r = p.add_run()
        r.text = ("• " if lvl == 0 else "– ") + txt
        r.font.size = Pt(21 if lvl == 0 else 17)
        r.font.color.rgb = DARK if lvl == 0 else GREY
        p.space_after = Pt(7)


def table_slide(title, headers, rows, col_widths=None):
    s = prs.slides.add_slide(BLANK)
    header(s, title)
    nrows, ncols = len(rows) + 1, len(headers)
    left, top = Inches(0.6), Inches(1.35)
    width = SW - Inches(1.2)
    height = Inches(5.6)
    gtbl = s.shapes.add_table(nrows, ncols, left, top, width, height)
    tbl = gtbl.table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)
    # header row
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = FABRIC
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; r = p.add_run(); r.text = h
        r.font.bold = True; r.font.size = Pt(15); r.font.color.rgb = WHITE
    # body rows
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 else FABRIC_LT
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; r = p.add_run(); r.text = val
            r.font.size = Pt(13); r.font.color.rgb = DARK
            if ci == 0:
                r.font.bold = True


def image_slide(title, caption, images):
    s = prs.slides.add_slide(BLANK)
    header(s, title)
    textbox(s, Inches(0.6), Inches(1.12), Inches(12.1), Inches(0.8), caption,
            size=15, color=GREY)
    paths = [SHOTS / n for n in images if (SHOTS / n).exists()]
    if not paths:
        textbox(s, Inches(0.6), Inches(3), Inches(12), Inches(1),
                "[screenshot missing: " + ", ".join(images) + "]", size=16, color=GREY)
        return
    from PIL import Image
    area_top, area_h, gap = Inches(2.0), Inches(4.9), Inches(0.3)
    n = len(paths)
    each_w = int((SW - Inches(1.2) - gap * (n - 1)) / n)
    x = Inches(0.6)
    for p in paths:
        with Image.open(p) as im:
            ar = im.height / im.width
        w = each_w; h = int(w * ar)
        if h > area_h:
            h = int(area_h); w = int(h / ar)
        y = area_top + int((area_h - h) / 2)
        s.shapes.add_picture(str(p), x, y, width=w, height=h)
        x += each_w + gap


# ---- Build ------------------------------------------------------------------
title_slide(
    "Fabric-as-Code",
    "Deploy Microsoft Fabric end-to-end — from nothing to a running platform",
    "Azure CLI  ·  Bicep  ·  Fabric REST APIs",
)

table_slide(
    "The problem: two control planes",
    ["Layer", "What it manages", "Tooling here"],
    [
        ["Azure", "Fabric capacity (Microsoft.Fabric/capacities) — billable compute", "Bicep + az"],
        ["Fabric", "Workspaces + items (Lakehouse, Warehouse, Notebook, Pipeline)", "Fabric REST"],
        ["Data", "Objects inside an item — e.g. stored procedures", "T-SQL"],
    ],
    col_widths=[1.6, 7.7, 2.4],
)

table_slide(
    "The building blocks — what each part does",
    ["Component", "What it is / does"],
    [
        ["Capacity (F-SKU)", "The compute every item runs on. Billed hourly; pause to stop."],
        ["Workspace", "Logical container for items; must sit on a capacity."],
        ["Lakehouse", "OneLake files + Delta tables; auto SQL endpoint (read-only T-SQL)."],
        ["Warehouse", "Full read/write T-SQL engine: tables, views, stored procs, transactions."],
        ["Notebook", "Spark/PySpark compute; bound to a Lakehouse to read/write tables."],
        ["Data Pipeline", "Orchestrator chaining notebooks, copies and stored procedures."],
    ],
    col_widths=[2.6, 9.1],
)

content_slide("What the automation does (6 steps)", [
    ("Provision resource group + Fabric capacity — 02 → infra/capacity.bicep", 0),
    ("Create the workspace — 03 → POST /v1/workspaces", 0),
    ("Assign workspace to capacity — 04 → assignToCapacity (unlocks items)", 0),
    ("Deploy items — 05 → Lakehouse, Warehouse, Notebook, Pipeline", 0),
    ("Deploy stored procedures — 06 → T-SQL to the Warehouse SQL endpoint", 0),
    ("Optional: connect to Git — 07 → workspace under source control", 0),
    ("Run individually to learn, or all at once with deploy-all.", 0),
])

table_slide(
    "The deployment scripts (PowerShell + bash)",
    ["Script", "What it does"],
    [
        ["00-prerequisites", "Checks az, pwsh/bash, sqlcmd, and login state"],
        ["01-login", "az login — interactive or service principal"],
        ["02-provision-capacity", "Bicep deploy of the Fabric capacity"],
        ["03 / 04", "Create workspace · assign it to the capacity"],
        ["05-deploy-items", "Create the four items from definition files"],
        ["06-deploy-stored-procedures", "Run T-SQL against the Warehouse (Entra token)"],
        ["99-teardown", "Delete the workspace and the resource group"],
    ],
    col_widths=[3.2, 8.5],
)

content_slide("How it stays reliable", [
    ("Idempotent — existing workspaces/items are detected and reused, not duplicated", 0),
    ("State — resolved GUIDs saved to .state.json between steps", 0),
    ("Auth — one az session mints tokens for Fabric and SQL; no extra tooling for procs", 0),
    ("Long-running ops — REST helper polls Operation-Location until Succeeded", 0),
    ("Templated definitions — __TOKENS__ in notebook/pipeline JSON become real GUIDs", 0),
    ("Portable — same definitions deploy into any workspace/tenant", 0),
])

content_slide("Repeatable across tenants & RGs", [
    ("All tenant / subscription / secret values live in a git-ignored .env", 0),
    ("Duplicate per environment: .env.dev, .env.customerA, …", 0),
    ("Same scripts, different config → identical environment anywhere", 0),
    ("Service-principal auth for CI/CD and multi-tenant automation", 0),
    ("Public repo: contains no secrets", 0),
])

image_slide("Fabric — workspace contents",
            "Every item created via the Fabric REST API (Lakehouse + SQL endpoint, Warehouse, Notebook, Pipeline). Created by scripts 03–05.",
            ["01-workspace-list.png"])
image_slide("Fabric — Lakehouse with data",
            "orders_bronze Delta table (5 rows), written by the notebook run — the bronze landing layer. The auto SQL endpoint enables read-only T-SQL / BI.",
            ["02-lakehouse-table.png"])
image_slide("Fabric — Notebook",
            "nb_demo_load: PySpark builds the dataset and saveAsTable writes it to the attached Lakehouse. The Lakehouse binding is injected at deploy time.",
            ["03-notebook.png"])
image_slide("Fabric — Data Pipeline",
            "Design: Wait → Run Notebook (success dependency). Run: Succeeded. Orchestration that executed the notebook and populated the Delta table.",
            ["04-pipeline-canvas.png", "05-pipeline-run.png"])
image_slide("Azure — Fabric capacity",
            "The billable Microsoft.Fabric/capacities resource (F2) — the compute all items run on. Deployed by Bicep; pause to stop hourly billing.",
            ["06-azure-capacity.png"])

content_slide("What you can build on this", [
    ("Medallion platform — Lakehouse bronze/silver/gold + Spark notebooks", 0),
    ("SQL warehousing — Warehouse with versioned schema, tables & stored procedures", 0),
    ("Orchestration — Pipelines chaining notebooks, copies & procedures", 0),
    ("Promotion — Git integration + deployment pipelines (dev → test → prod)", 0),
    ("Governance as code — capacity sizing, admins & RBAC in source control", 0),
])

content_slide("Get started", [
    ("git clone https://github.com/<owner>/fabric-as-code", 0),
    ("cp .env.example .env   (fill in tenant / subscription / capacity)", 0),
    ("pwsh ./scripts/powershell/deploy-all.ps1", 0),
    ("…or  ./scripts/bash/deploy-all.sh", 0),
    ("Tear it all down with 99-teardown.", 0),
])

prs.save(OUT)
print(f"Wrote {OUT}")
